"""
Unified Content Service - Consolidated content management
Combines all content-related services into a single, optimized service
"""

import asyncio
import logging
from typing import Dict, List, Optional, Any, Union, BinaryIO
from dataclasses import dataclass
from enum import Enum
import json
import uuid
from datetime import datetime
import aiofiles
import aiohttp
from PIL import Image
import cv2
import numpy as np
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from openpyxl import Workbook
import markdown
from jinja2 import Template
import hashlib
import mimetypes

logger = logging.getLogger(__name__)

class ContentType(Enum):
    """Content Types"""
    PRESENTATION = "presentation"
    DOCUMENT = "document"
    WEB_PAGE = "web_page"
    IMAGE = "image"
    VIDEO = "video"
    AUDIO = "audio"
    PDF = "pdf"
    EXCEL = "excel"
    JSON = "json"
    MARKDOWN = "markdown"

class ContentStatus(Enum):
    """Content Status"""
    DRAFT = "draft"
    PUBLISHED = "published"
    ARCHIVED = "archived"
    DELETED = "deleted"

@dataclass
class ContentMetadata:
    """Content Metadata"""
    id: str
    title: str
    description: str
    content_type: ContentType
    status: ContentStatus
    created_at: datetime
    updated_at: datetime
    author: str
    tags: List[str]
    size: int
    format: str
    version: int
    parent_id: Optional[str] = None
    thumbnail_url: Optional[str] = None
    custom_fields: Dict[str, Any] = None

@dataclass
class ContentRequest:
    """Content Generation Request"""
    content_type: ContentType
    title: str
    description: str
    template: Optional[str] = None
    data: Optional[Dict[str, Any]] = None
    style: Optional[Dict[str, Any]] = None
    author: str = "system"
    tags: List[str] = None

@dataclass
class ContentResponse:
    """Content Generation Response"""
    content_id: str
    content_url: str
    metadata: ContentMetadata
    generation_time: float
    success: bool
    error: Optional[str] = None

class UnifiedContentService:
    """
    Unified Content Service - Consolidated content management
    Handles all content types: presentations, documents, web pages, media, etc.
    """
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.storage_path = config.get("storage_path", "./content_storage")
        self.templates_path = config.get("templates_path", "./templates")
        self.thumbnails_path = config.get("thumbnails_path", "./thumbnails")
        
        # Content registry
        self.content_registry: Dict[str, ContentMetadata] = {}
        
        # Templates
        self.templates: Dict[str, Template] = {}
        
        logger.info("UnifiedContentService initialized")
    
    async def generate_content(self, request: ContentRequest) -> ContentResponse:
        """Generate content based on request"""
        start_time = asyncio.get_event_loop().time()
        content_id = str(uuid.uuid4())
        
        try:
            if request.content_type == ContentType.PRESENTATION:
                content_url = await self._generate_presentation(content_id, request)
            elif request.content_type == ContentType.DOCUMENT:
                content_url = await self._generate_document(content_id, request)
            elif request.content_type == ContentType.WEB_PAGE:
                content_url = await self._generate_web_page(content_id, request)
            elif request.content_type == ContentType.PDF:
                content_url = await self._generate_pdf(content_id, request)
            elif request.content_type == ContentType.EXCEL:
                content_url = await self._generate_excel(content_id, request)
            else:
                raise ValueError(f"Unsupported content type: {request.content_type}")
            
            # Create metadata
            metadata = ContentMetadata(
                id=content_id,
                title=request.title,
                description=request.description,
                content_type=request.content_type,
                status=ContentStatus.DRAFT,
                created_at=datetime.now(),
                updated_at=datetime.now(),
                author=request.author,
                tags=request.tags or [],
                size=await self._get_file_size(content_url),
                format=request.content_type.value,
                version=1
            )
            
            # Register content
            self.content_registry[content_id] = metadata
            
            generation_time = asyncio.get_event_loop().time() - start_time
            
            return ContentResponse(
                content_id=content_id,
                content_url=content_url,
                metadata=metadata,
                generation_time=generation_time,
                success=True
            )
            
        except Exception as e:
            logger.error(f"Error generating content: {e}")
            return ContentResponse(
                content_id=content_id,
                content_url="",
                metadata=None,
                generation_time=asyncio.get_event_loop().time() - start_time,
                success=False,
                error=str(e)
            )
    
    async def _generate_presentation(self, content_id: str, request: ContentRequest) -> str:
        """Generate PowerPoint presentation"""
        try:
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = request.title
            subtitle.text = request.description
            
            # Add content slides based on data
            if request.data and "slides" in request.data:
                for slide_data in request.data["slides"]:
                    content_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(content_slide_layout)
                    
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_data.get("title", "")
                    content.text = slide_data.get("content", "")
            
            # Save presentation
            file_path = f"{self.storage_path}/presentation_{content_id}.pptx"
            prs.save(file_path)
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise
    
    async def _generate_document(self, content_id: str, request: ContentRequest) -> str:
        """Generate document (Markdown/HTML)"""
        try:
            # Create markdown content
            markdown_content = f"# {request.title}\n\n{request.description}\n\n"
            
            if request.data:
                for section in request.data.get("sections", []):
                    markdown_content += f"## {section.get('title', '')}\n\n{section.get('content', '')}\n\n"
            
            # Convert to HTML
            html_content = markdown.markdown(markdown_content)
            
            # Apply template if provided
            if request.template:
                template = await self._load_template(request.template)
                html_content = template.render(
                    title=request.title,
                    content=html_content,
                    style=request.style or {}
                )
            
            # Save document
            file_path = f"{self.storage_path}/document_{content_id}.html"
            async with aiofiles.open(file_path, 'w', encoding='utf-8') as f:
                await f.write(html_content)
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error generating document: {e}")
            raise
    
    async def _generate_web_page(self, content_id: str, request: ContentRequest) -> str:
        """Generate responsive web page"""
        try:
            # HTML template
            html_template = """
            <!DOCTYPE html>
            <html lang="en">
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
                <title>{{ title }}</title>
                <style>
                    body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
                    .container { max-width: 1200px; margin: 0 auto; }
                    .header { text-align: center; margin-bottom: 40px; }
                    .content { line-height: 1.6; }
                    .responsive { width: 100%; height: auto; }
                    @media (max-width: 768px) {
                        .container { padding: 10px; }
                    }
                </style>
            </head>
            <body>
                <div class="container">
                    <div class="header">
                        <h1>{{ title }}</h1>
                        <p>{{ description }}</p>
                    </div>
                    <div class="content">
                        {{ content }}
                    </div>
                </div>
            </body>
            </html>
            """
            
            template = Template(html_template)
            
            # Generate content
            content_html = ""
            if request.data:
                for section in request.data.get("sections", []):
                    content_html += f"<section><h2>{section.get('title', '')}</h2><p>{section.get('content', '')}</p></section>"
            
            # Render HTML
            html_content = template.render(
                title=request.title,
                description=request.description,
                content=content_html
            )
            
            # Save web page
            file_path = f"{self.storage_path}/webpage_{content_id}.html"
            async with aiofiles.open(file_path, 'w', encoding='utf-8') as f:
                await f.write(html_content)
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error generating web page: {e}")
            raise
    
    async def _generate_pdf(self, content_id: str, request: ContentRequest) -> str:
        """Generate PDF document"""
        try:
            file_path = f"{self.storage_path}/document_{content_id}.pdf"
            doc = SimpleDocTemplate(file_path, pagesize=letter)
            
            # Styles
            styles = getSampleStyleSheet()
            title_style = styles['Title']
            heading_style = styles['Heading1']
            normal_style = styles['Normal']
            
            # Content
            story = []
            story.append(Paragraph(request.title, title_style))
            story.append(Spacer(1, 12))
            story.append(Paragraph(request.description, normal_style))
            story.append(Spacer(1, 12))
            
            if request.data:
                for section in request.data.get("sections", []):
                    story.append(Paragraph(section.get("title", ""), heading_style))
                    story.append(Paragraph(section.get("content", ""), normal_style))
                    story.append(Spacer(1, 12))
            
            # Build PDF
            doc.build(story)
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error generating PDF: {e}")
            raise
    
    async def _generate_excel(self, content_id: str, request: ContentRequest) -> str:
        """Generate Excel spreadsheet"""
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = request.title[:31]  # Excel sheet name limit
            
            # Add title
            ws['A1'] = request.title
            ws['A2'] = request.description
            
            # Add data
            if request.data and "data" in request.data:
                data = request.data["data"]
                if isinstance(data, list) and len(data) > 0:
                    # Headers
                    headers = list(data[0].keys())
                    for col, header in enumerate(headers, 1):
                        ws.cell(row=3, column=col, value=header)
                    
                    # Data rows
                    for row, item in enumerate(data, 4):
                        for col, header in enumerate(headers, 1):
                            ws.cell(row=row, column=col, value=item.get(header, ""))
            
            # Save Excel file
            file_path = f"{self.storage_path}/spreadsheet_{content_id}.xlsx"
            wb.save(file_path)
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error generating Excel: {e}")
            raise
    
    async def _load_template(self, template_name: str) -> Template:
        """Load Jinja2 template"""
        if template_name not in self.templates:
            template_path = f"{self.templates_path}/{template_name}.html"
            async with aiofiles.open(template_path, 'r', encoding='utf-8') as f:
                template_content = await f.read()
                self.templates[template_name] = Template(template_content)
        
        return self.templates[template_name]
    
    async def _get_file_size(self, file_path: str) -> int:
        """Get file size in bytes"""
        try:
            async with aiofiles.open(file_path, 'rb') as f:
                content = await f.read()
                return len(content)
        except:
            return 0
    
    async def get_content(self, content_id: str) -> Optional[ContentMetadata]:
        """Get content metadata by ID"""
        return self.content_registry.get(content_id)
    
    async def list_content(self, 
                          content_type: Optional[ContentType] = None,
                          status: Optional[ContentStatus] = None,
                          author: Optional[str] = None,
                          limit: int = 100) -> List[ContentMetadata]:
        """List content with filters"""
        content_list = list(self.content_registry.values())
        
        # Apply filters
        if content_type:
            content_list = [c for c in content_list if c.content_type == content_type]
        if status:
            content_list = [c for c in content_list if c.status == status]
        if author:
            content_list = [c for c in content_list if c.author == author]
        
        # Sort by creation date (newest first)
        content_list.sort(key=lambda x: x.created_at, reverse=True)
        
        return content_list[:limit]
    
    async def update_content(self, content_id: str, updates: Dict[str, Any]) -> bool:
        """Update content metadata"""
        try:
            if content_id in self.content_registry:
                metadata = self.content_registry[content_id]
                
                # Update fields
                for key, value in updates.items():
                    if hasattr(metadata, key):
                        setattr(metadata, key, value)
                
                metadata.updated_at = datetime.now()
                metadata.version += 1
                
                return True
            
            return False
            
        except Exception as e:
            logger.error(f"Error updating content {content_id}: {e}")
            return False
    
    async def delete_content(self, content_id: str) -> bool:
        """Delete content"""
        try:
            if content_id in self.content_registry:
                metadata = self.content_registry[content_id]
                metadata.status = ContentStatus.DELETED
                return True
            
            return False
            
        except Exception as e:
            logger.error(f"Error deleting content {content_id}: {e}")
            return False
    
    async def search_content(self, query: str, content_type: Optional[ContentType] = None) -> List[ContentMetadata]:
        """Search content by query"""
        try:
            results = []
            query_lower = query.lower()
            
            for metadata in self.content_registry.values():
                if metadata.status == ContentStatus.DELETED:
                    continue
                
                if content_type and metadata.content_type != content_type:
                    continue
                
                # Search in title, description, and tags
                if (query_lower in metadata.title.lower() or
                    query_lower in metadata.description.lower() or
                    any(query_lower in tag.lower() for tag in metadata.tags)):
                    results.append(metadata)
            
            # Sort by relevance (title matches first)
            results.sort(key=lambda x: (
                query_lower not in x.title.lower(),
                x.created_at
            ))
            
            return results
            
        except Exception as e:
            logger.error(f"Error searching content: {e}")
            return []
    
    async def generate_thumbnail(self, content_id: str, content_url: str) -> Optional[str]:
        """Generate thumbnail for content"""
        try:
            # Determine content type and generate appropriate thumbnail
            metadata = self.content_registry.get(content_id)
            if not metadata:
                return None
            
            if metadata.content_type == ContentType.IMAGE:
                # For images, create a resized version
                thumbnail_path = f"{self.thumbnails_path}/thumb_{content_id}.jpg"
                await self._create_image_thumbnail(content_url, thumbnail_path)
                return thumbnail_path
            
            elif metadata.content_type == ContentType.WEB_PAGE:
                # For web pages, create a screenshot (placeholder)
                thumbnail_path = f"{self.thumbnails_path}/thumb_{content_id}.jpg"
                await self._create_web_thumbnail(content_url, thumbnail_path)
                return thumbnail_path
            
            # Default thumbnail
            return None
            
        except Exception as e:
            logger.error(f"Error generating thumbnail for {content_id}: {e}")
            return None
    
    async def _create_image_thumbnail(self, image_path: str, thumbnail_path: str, size: tuple = (300, 200)):
        """Create image thumbnail"""
        try:
            image = Image.open(image_path)
            image.thumbnail(size, Image.Resampling.LANCZOS)
            image.save(thumbnail_path, "JPEG", quality=85)
        except Exception as e:
            logger.error(f"Error creating image thumbnail: {e}")
    
    async def _create_web_thumbnail(self, web_path: str, thumbnail_path: str):
        """Create web page thumbnail (placeholder)"""
        try:
            # In a real implementation, this would use a headless browser
            # to take a screenshot of the web page
            pass
        except Exception as e:
            logger.error(f"Error creating web thumbnail: {e}")
    
    async def export_content(self, content_id: str, export_format: str) -> Optional[str]:
        """Export content in different formats"""
        try:
            metadata = self.content_registry.get(content_id)
            if not metadata:
                return None
            
            # Get original content
            original_path = f"{self.storage_path}/{metadata.content_type.value}_{content_id}.{metadata.format}"
            
            if export_format == "json":
                export_path = f"{self.storage_path}/export_{content_id}.json"
                export_data = {
                    "metadata": {
                        "id": metadata.id,
                        "title": metadata.title,
                        "description": metadata.description,
                        "content_type": metadata.content_type.value,
                        "created_at": metadata.created_at.isoformat(),
                        "author": metadata.author,
                        "tags": metadata.tags
                    }
                }
                async with aiofiles.open(export_path, 'w', encoding='utf-8') as f:
                    await f.write(json.dumps(export_data, indent=2))
                return export_path
            
            # Add more export formats as needed
            return original_path
            
        except Exception as e:
            logger.error(f"Error exporting content {content_id}: {e}")
            return None
    
    async def get_content_statistics(self) -> Dict[str, Any]:
        """Get content statistics"""
        try:
            total_content = len(self.content_registry)
            content_by_type = {}
            content_by_status = {}
            
            for metadata in self.content_registry.values():
                # Count by type
                content_type = metadata.content_type.value
                content_by_type[content_type] = content_by_type.get(content_type, 0) + 1
                
                # Count by status
                status = metadata.status.value
                content_by_status[status] = content_by_status.get(status, 0) + 1
            
            return {
                "total_content": total_content,
                "content_by_type": content_by_type,
                "content_by_status": content_by_status,
                "storage_used": sum(metadata.size for metadata in self.content_registry.values())
            }
            
        except Exception as e:
            logger.error(f"Error getting content statistics: {e}")
            return {}
    
    async def health_check(self) -> Dict[str, Any]:
        """Health check for content service"""
        try:
            stats = await self.get_content_statistics()
            
            return {
                "status": "healthy",
                "total_content": stats.get("total_content", 0),
                "storage_path": self.storage_path,
                "templates_loaded": len(self.templates)
            }
            
        except Exception as e:
            logger.error(f"Health check failed: {e}")
            return {"status": "unhealthy", "error": str(e)}


























