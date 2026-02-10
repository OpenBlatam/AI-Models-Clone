"""PowerPoint Converter - Convert Markdown to PowerPoint format"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any, Optional

from .base_converter import BaseConverter


class PPTConverter(BaseConverter):
    """Convert Markdown to PowerPoint (.pptx) format"""
    
    async def convert(
        self,
        parsed_content: Dict[str, Any],
        output_path: str,
        include_charts: bool = True,
        include_tables: bool = True,
        custom_styling: Optional[Dict[str, Any]] = None
    ) -> None:
        """Convert to PowerPoint format"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        if parsed_content.get("title"):
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = parsed_content["title"]
            subtitle.text = "Generated from Markdown"
        
        # Content slides
        current_slide = None
        
        for heading in parsed_content.get("headings", []):
            if heading["level"] == 1:
                # New slide for H1
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = heading["text"]
                tf = body_shape.text_frame
                tf.text = ""
            elif current_slide:
                # Add as bullet point
                body_shape = current_slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                p = tf.add_paragraph()
                p.text = heading["text"]
                p.level = heading["level"] - 1
                p.font.size = Pt(18 - heading["level"] * 2)
        
        # Add table slides
        if include_tables:
            for table_data in parsed_content.get("tables", []):
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add title
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(0.8)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = "Table"
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.bold = True
                
                # Add table
                rows = len(table_data["rows"]) + 1
                cols = len(table_data["headers"])
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Headers
                for col_idx, header in enumerate(table_data["headers"]):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                # Data rows
                for row_idx, row_data in enumerate(table_data["rows"], start=1):
                    for col_idx, value in enumerate(row_data):
                        if col_idx < cols:
                            table.cell(row_idx, col_idx).text = str(value)
        
        # Save presentation
        prs.save(output_path)

